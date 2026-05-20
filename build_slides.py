from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x21, 0x7A, 0x3C)
TBLHDR = RGBColor(0x1F, 0x38, 0x64)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_text(slide, text, x, y, w, h,
             size=16, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_rect(slide, x, y, w, h, fill=LGRAY, line=None):
    from pptx.util import Pt as UPt
    shp = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp


def slide_heading(slide, text, y=Inches(0.38)):
    add_text(slide, text,
             Inches(0.6), y, Inches(12.0), Inches(0.6),
             size=26, bold=True, color=NAVY)


def divider(slide, y):
    """Thin navy line under heading."""
    ln = slide.shapes.add_shape(1,
         Inches(0.6), y, Inches(12.1), Inches(0.03))
    ln.fill.solid()
    ln.fill.fore_color.rgb = NAVY
    ln.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)

# Navy block on left 40%
add_rect(s1, 0, 0, Inches(5.3), H, fill=NAVY)

# Title on navy block
add_text(s1, "Bitcoin Price\nForecasting",
         Inches(0.45), Inches(2.0), Inches(4.4), Inches(2.2),
         size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s1, "ML Semester Project  ·  Mid Evaluation",
         Inches(0.45), Inches(4.35), Inches(4.4), Inches(0.5),
         size=14, color=RGBColor(0xB0, 0xBE, 0xD4), align=PP_ALIGN.LEFT)

add_text(s1, "Esha Saeed  ·  Khadija Masood  ·  Muzammil Mahmood",
         Inches(0.45), Inches(5.0), Inches(4.6), Inches(0.45),
         size=12, color=RGBColor(0xB0, 0xBE, 0xD4))

# Right side — course info
add_text(s1, "CS373 Machine Learning",
         Inches(5.8), Inches(3.1), Inches(6.8), Inches(0.5),
         size=20, bold=True, color=NAVY)
add_text(s1, "Information Technology University",
         Inches(5.8), Inches(3.65), Inches(6.8), Inches(0.5),
         size=15, color=GRAY)
add_text(s1, "2026",
         Inches(5.8), Inches(4.15), Inches(6.8), Inches(0.4),
         size=13, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)

add_text(s2,
         "Can machine learning predict whether\nBitcoin's price will rise or fall tomorrow?",
         Inches(0.6), Inches(0.3), Inches(12.1), Inches(1.3),
         size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

divider(s2, Inches(1.65))

# Left box
add_rect(s2, Inches(0.6), Inches(1.85), Inches(5.7), Inches(4.8), fill=LGRAY)
add_text(s2, "Why it's hard",
         Inches(0.85), Inches(2.0), Inches(5.2), Inches(0.45),
         size=15, bold=True, color=NAVY)
for i, line in enumerate([
    "Price is non-stationary — mean and variance drift over time",
    "Volatile and non-linear — no simple formula explains moves",
    "No single indicator predicts it reliably on its own",
]):
    add_text(s2, f"–   {line}",
             Inches(0.95), Inches(2.55) + Inches(0.75)*i,
             Inches(5.0), Inches(0.65),
             size=13, color=BLACK)

# Right box
add_rect(s2, Inches(7.0), Inches(1.85), Inches(5.7), Inches(4.8), fill=LGRAY)
add_text(s2, "Our approach",
         Inches(7.25), Inches(2.0), Inches(5.2), Inches(0.45),
         size=15, bold=True, color=NAVY)
for i, line in enumerate([
    "4,499 days of BTC/USD data  (2014 – 2026)",
    "10 engineered technical indicators as features",
    "5 ML models benchmarked against a persistence baseline",
]):
    add_text(s2, f"–   {line}",
             Inches(7.35), Inches(2.55) + Inches(0.75)*i,
             Inches(5.1), Inches(0.65),
             size=13, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
slide_heading(s3, "Data Pipeline")
divider(s3, Inches(1.05))

steps = [
    "Raw CSV\n7.5M rows\n1-min bars",
    "Daily\nAggregation\n4,499 days",
    "ADF\nStationarity\nTest",
    "Log-\nDifferencing",
    "Feature\nMatrix\n4,485 × 10",
]
box_w = Inches(2.0)
box_h = Inches(1.55)
gap   = Inches(0.38)
start_x = Inches(0.6)
y_row = Inches(1.5)

for i, label in enumerate(steps):
    bx = start_x + (box_w + gap) * i
    add_rect(s3, bx, y_row, box_w, box_h,
             fill=NAVY if i in (0, 4) else LGRAY)
    fc = WHITE if i in (0, 4) else NAVY
    add_text(s3, label, bx + Inches(0.1), y_row + Inches(0.15),
             box_w - Inches(0.2), box_h - Inches(0.2),
             size=12, bold=True, color=fc, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = bx + box_w + Inches(0.08)
        add_text(s3, "→", ax, y_row + Inches(0.55),
                 Inches(0.22), Inches(0.45),
                 size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Finding box
add_rect(s3, Inches(0.6), Inches(3.45), Inches(12.1), Inches(1.25),
         fill=RGBColor(0xE8, 0xEC, 0xF4), line=NAVY)
add_text(s3,
         "ADF test confirmed raw prices are non-stationary (p = 0.78).\n"
         "Log returns are stationary (p ≈ 0) — used as the model target across all models.",
         Inches(0.9), Inches(3.58), Inches(11.5), Inches(1.0),
         size=13, color=BLACK, italic=True)

# Note on trimming
add_text(s3,
         "Data trimmed to most recent 4,500 days — early illiquid BTC era excluded to improve signal quality.",
         Inches(0.6), Inches(5.05), Inches(12.1), Inches(0.45),
         size=11, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Features
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
slide_heading(s4, "What We Feed the Models")
divider(s4, Inches(1.05))

# Left column
add_rect(s4, Inches(0.6), Inches(1.3), Inches(5.7), Inches(5.4), fill=LGRAY)
add_text(s4, "10 Input Features",
         Inches(0.85), Inches(1.45), Inches(5.2), Inches(0.4),
         size=15, bold=True, color=NAVY)

features = [
    ("Close, Volume", "raw market data"),
    ("EMA-9,  EMA-21", "short & long trend"),
    ("MACD,  Signal", "momentum divergence"),
    ("RSI-14", "overbought / oversold"),
    ("Momentum,  PROC", "price rate of change"),
    ("Stochastic %K", "position in recent range"),
]
for i, (feat, desc) in enumerate(features):
    y = Inches(2.0) + Inches(0.72) * i
    add_text(s4, feat,  Inches(0.9),  y, Inches(2.6), Inches(0.38),
             size=13, bold=True, color=BLACK)
    add_text(s4, desc, Inches(3.55), y, Inches(2.6), Inches(0.38),
             size=12, color=GRAY, italic=True)

# Right column
add_rect(s4, Inches(7.0), Inches(1.3), Inches(5.7), Inches(2.45), fill=LGRAY)
add_text(s4, "Regression Target",
         Inches(7.25), Inches(1.45), Inches(5.2), Inches(0.4),
         size=15, bold=True, color=NAVY)
add_text(s4, "Next-day log return  (continuous)",
         Inches(7.25), Inches(1.95), Inches(5.2), Inches(0.35),
         size=13, color=BLACK)
add_text(s4, "log(Closeₜ₊₁) − log(Closeₜ)",
         Inches(7.25), Inches(2.35), Inches(5.2), Inches(0.45),
         size=12, color=GRAY, italic=True)

add_rect(s4, Inches(7.0), Inches(3.95), Inches(5.7), Inches(2.45), fill=LGRAY)
add_text(s4, "Classification Target",
         Inches(7.25), Inches(4.1), Inches(5.2), Inches(0.4),
         size=15, bold=True, color=NAVY)
add_text(s4, "Direction:  UP (1)  or  DOWN (0)",
         Inches(7.25), Inches(4.6), Inches(5.2), Inches(0.35),
         size=13, color=BLACK)
add_text(s4, "52.1% up-days  —  nearly balanced classes",
         Inches(7.25), Inches(5.0), Inches(5.2), Inches(0.35),
         size=12, color=GRAY, italic=True)

add_text(s4,
         "All indicators are backward-looking only.  "
         "Train/test split is time-ordered (80/20).  No data leakage.",
         Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.4),
         size=11, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Validation Strategy
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
slide_heading(s5, "Validation Strategy")
divider(s5, Inches(1.05))

blocks = [
    ("80 / 20 Temporal Split",
     "Train:  Jan 2014 – Nov 2023   (3,588 days)\nTest:    Nov 2023 – May 2026   (897 days)"),
    ("No Shuffling — Ever",
     "Shuffling mixes future data into training.\nAll splits strictly respect time order."),
    ("5-Fold TimeSeriesSplit  (hyperparameter tuning)",
     "Each fold trains on the past and validates on the immediate future.\nStandard K-Fold CV causes look-ahead bias on time series — not used."),
]

for i, (title, body) in enumerate(blocks):
    y = Inches(1.4) + Inches(1.75) * i
    add_rect(s5, Inches(0.6), y, Inches(12.1), Inches(1.55), fill=LGRAY)
    add_text(s5, title,
             Inches(0.9), y + Inches(0.12), Inches(11.5), Inches(0.45),
             size=14, bold=True, color=NAVY)
    add_text(s5, body,
             Inches(0.9), y + Inches(0.6), Inches(11.5), Inches(0.85),
             size=13, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Regression Results
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
slide_heading(s6, "Regression Results  —  Predicting Next-Day Return")
divider(s6, Inches(1.05))

# Table
col_x  = [Inches(0.6), Inches(5.5), Inches(8.0), Inches(10.0)]
col_w  = [Inches(4.7), Inches(2.3), Inches(1.8),  Inches(2.8)]
row_h  = Inches(0.6)
hdr_y  = Inches(1.3)
headers = ["Model", "RMSE", "MAE", "vs Baseline"]

# Header row
for j, hdr in enumerate(headers):
    add_rect(s6, col_x[j], hdr_y, col_w[j], row_h, fill=TBLHDR)
    add_text(s6, hdr, col_x[j] + Inches(0.1), hdr_y + Inches(0.12),
             col_w[j] - Inches(0.15), row_h - Inches(0.1),
             size=13, bold=True, color=WHITE)

rows = [
    ("Persistence Baseline", "0.0363", "0.0266", "—"),
    ("Linear Regression",    "0.0254", "0.0183", "✓  beats baseline"),
    ("Ridge  (λ = 10)",       "0.0251", "0.0180", "✓  beats baseline"),
    ("Lasso  (λ = 0.01)",     "0.0249", "0.0178", "✓  beats baseline"),
    ("Random Forest",        "0.0270", "0.0193", "✓  beats baseline"),
]

for i, row in enumerate(rows):
    ry = hdr_y + row_h * (i + 1)
    bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else LGRAY
    for j, cell in enumerate(row):
        add_rect(s6, col_x[j], ry, col_w[j], row_h, fill=bg)
        fc = GREEN if "✓" in cell else (NAVY if i == 0 else BLACK)
        add_text(s6, cell,
                 col_x[j] + Inches(0.1), ry + Inches(0.12),
                 col_w[j] - Inches(0.15), row_h - Inches(0.1),
                 size=13, bold=(i == 0), color=fc)

add_text(s6,
         "Note: Negative R² is expected for financial return prediction  "
         "— BTC return variance is inherently large.  RMSE vs baseline is the meaningful measure.",
         Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.5),
         size=11, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Classification Results
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
slide_heading(s7, "Classification Results  —  Predicting Direction (UP / DOWN)")
divider(s7, Inches(1.05))

col_x2 = [Inches(0.6), Inches(4.6), Inches(6.8), Inches(9.0), Inches(11.2)]
col_w2 = [Inches(3.8), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
hdr_y2 = Inches(1.3)
headers2 = ["Model", "Accuracy", "Precision", "Recall", "F1"]

for j, hdr in enumerate(headers2):
    add_rect(s7, col_x2[j], hdr_y2, col_w2[j], row_h, fill=TBLHDR)
    add_text(s7, hdr, col_x2[j] + Inches(0.1), hdr_y2 + Inches(0.12),
             col_w2[j] - Inches(0.15), row_h - Inches(0.1),
             size=13, bold=True, color=WHITE)

rows2 = [
    ("Majority Baseline",    "51.2%", "51.2%", "100%",  "67.7%"),
    ("Logistic Regression",  "52.4%", "55.4%", "35.7%", "43.4%"),
    ("Random Forest",        "50.2%", "55.5%", "13.3%", "21.4%"),
]
for i, row in enumerate(rows2):
    ry = hdr_y2 + row_h * (i + 1)
    bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else LGRAY
    for j, cell in enumerate(row):
        add_rect(s7, col_x2[j], ry, col_w2[j], row_h, fill=bg)
        add_text(s7, cell,
                 col_x2[j] + Inches(0.1), ry + Inches(0.12),
                 col_w2[j] - Inches(0.15), row_h - Inches(0.1),
                 size=13, bold=(i == 0), color=BLACK)

# Insight box
add_rect(s7, Inches(0.6), Inches(3.95), Inches(12.1), Inches(1.6),
         fill=RGBColor(0xE8, 0xEC, 0xF4), line=NAVY)
add_text(s7, "Key Insight",
         Inches(0.9), Inches(4.07), Inches(11.5), Inches(0.38),
         size=13, bold=True, color=NAVY)
add_text(s7,
         "Both models achieve ~55% precision on UP days — they are right more than half the time "
         "when predicting a bullish move.\n"
         "Logistic Regression has higher recall (catches more UP days). "
         "Random Forest is more conservative — fewer but higher-confidence signals.",
         Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.95),
         size=12, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What's Next
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
slide_heading(s8, "Remaining Work")
divider(s8, Inches(1.05))

# Left
add_rect(s8, Inches(0.6), Inches(1.3), Inches(5.7), Inches(5.0), fill=LGRAY)
add_text(s8, "Final Model",
         Inches(0.85), Inches(1.45), Inches(5.2), Inches(0.4),
         size=15, bold=True, color=NAVY)
model_items = [
    "SVM with RBF kernel",
    "Applied to both regression and classification",
    "Kernel comparison: Linear vs Polynomial vs RBF",
    "C and γ tuned via TimeSeriesSplit GridSearchCV",
]
for i, item in enumerate(model_items):
    add_text(s8, f"–   {item}",
             Inches(0.9), Inches(2.05) + Inches(0.72)*i,
             Inches(5.1), Inches(0.62),
             size=13, color=BLACK)

# Right
add_rect(s8, Inches(7.0), Inches(1.3), Inches(5.7), Inches(5.0), fill=LGRAY)
add_text(s8, "Dashboard & API",
         Inches(7.25), Inches(1.45), Inches(5.2), Inches(0.4),
         size=15, bold=True, color=NAVY)
sys_items = [
    "Flask REST API — serving all model results",
    "React dashboard:",
    "    · Historical price + indicator overlays",
    "    · Model comparison (RMSE / Accuracy)",
    "    · Actual vs predicted chart",
    "    · Feature importance visualisation",
]
for i, item in enumerate(sys_items):
    add_text(s8, f"–   {item}" if not item.startswith("    ") else item,
             Inches(7.3), Inches(2.05) + Inches(0.65)*i,
             Inches(5.1), Inches(0.58),
             size=13, color=BLACK)


# ── Save ───────────────────────────────────────────────────────────────────────
out = r"D:\ML Proj\Mid_Eval_Slides.pptx"
prs.save(out)
print(f"Saved: {out}")
