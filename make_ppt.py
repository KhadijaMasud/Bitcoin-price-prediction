from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BG      = RGBColor(0xFA, 0xFA, 0xF8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x18)
MUTED   = RGBColor(0x6B, 0x6B, 0x68)
LIGHT   = RGBColor(0xE8, 0xE8, 0xE3)
ACCENT  = RGBColor(0x18, 0x5F, 0xA5)
ACCENT2 = RGBColor(0xE6, 0xF1, 0xFB)
GREEN   = RGBColor(0x3B, 0x6D, 0x11)
GREEN2  = RGBColor(0xEA, 0xF3, 0xDE)
AMBER   = RGBColor(0x85, 0x4F, 0x0B)
AMBER2  = RGBColor(0xFA, 0xEE, 0xDA)

BLANK = prs.slide_layouts[6]

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl

def box(sl, x, y, w, h, fill=WHITE, line=LIGHT, line_w=Pt(0.5)):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = line_w
    return s

def label(sl, text, x, y, w, h, size=11, bold=False, color=MUTED, align=PP_ALIGN.LEFT):
    tf = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tf

def slide_num(sl, n):
    label(sl, f"{n:02d} / 05", 0.4, 0.22, 1.5, 0.25, size=9, color=MUTED)

def title_bar(sl, text):
    box(sl, 0.4, 0.45, 12.53, 0.045, fill=DARK, line=DARK)
    label(sl, text, 0.4, 0.52, 12.0, 0.45, size=18, bold=False, color=DARK)

def tag(sl, text, x, y, fill=ACCENT2, tc=ACCENT):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(len(text)*0.085+0.2), Inches(0.22))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill
    tf = s.text_frame; tf.margin_left=Emu(60000); tf.margin_right=Emu(60000)
    tf.margin_top=Emu(20000); tf.margin_bottom=Emu(20000)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.color.rgb = tc; r.font.name = "Calibri"

def bullet_block(sl, items, x, y, w, color=DARK):
    cy = y
    for item in items:
        dot = sl.shapes.add_shape(9, Inches(x), Inches(cy+0.055), Inches(0.06), Inches(0.06))
        dot.fill.solid(); dot.fill.fore_color.rgb = MUTED
        dot.line.fill.background()
        tf = sl.shapes.add_textbox(Inches(x+0.14), Inches(cy), Inches(w-0.14), Inches(0.35))
        tf.text_frame.word_wrap = True
        p = tf.text_frame.paragraphs[0]; r = p.add_run(); r.text = item
        r.font.size = Pt(11); r.font.color.rgb = color; r.font.name = "Calibri"
        cy += 0.38
    return cy

def metric_card(sl, value, lbl, x, y, w=2.8, vc=DARK):
    s = box(sl, x, y, w, 0.75, fill=WHITE)
    tf = sl.shapes.add_textbox(Inches(x+0.15), Inches(y+0.06), Inches(w-0.3), Inches(0.35))
    p = tf.text_frame.paragraphs[0]; r = p.add_run(); r.text = value
    r.font.size = Pt(20); r.font.bold = False; r.font.color.rgb = vc; r.font.name = "Calibri"
    tf2 = sl.shapes.add_textbox(Inches(x+0.15), Inches(y+0.42), Inches(w-0.3), Inches(0.25))
    p2 = tf2.text_frame.paragraphs[0]; r2 = p2.add_run(); r2.text = lbl
    r2.font.size = Pt(9); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"

def info_card(sl, title, sub, x, y, w=3.9, h=0.75):
    box(sl, x, y, w, h, fill=WHITE)
    tf = sl.shapes.add_textbox(Inches(x+0.15), Inches(y+0.07), Inches(w-0.3), Inches(0.28))
    p = tf.text_frame.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(11); r.font.bold = False; r.font.color.rgb = DARK; r.font.name = "Calibri"
    tf2 = sl.shapes.add_textbox(Inches(x+0.15), Inches(y+0.35), Inches(w-0.3), Inches(h-0.4))
    tf2.text_frame.word_wrap = True
    p2 = tf2.text_frame.paragraphs[0]; r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(9.5); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"

def section_label(sl, text, x, y):
    label(sl, text.upper(), x, y, 6, 0.2, size=8.5, bold=True, color=MUTED)

def hbar(sl, model, acc_pct, x, y, highlight=False):
    tc = ACCENT if highlight else DARK
    label(sl, model, x, y, 3.6, 0.25, size=10, color=tc)
    bw = acc_pct * 0.055
    bf = ACCENT if highlight else LIGHT
    s = sl.shapes.add_shape(1, Inches(x+3.7), Inches(y+0.04), Inches(bw), Inches(0.14))
    s.fill.solid(); s.fill.fore_color.rgb = bf; s.line.fill.background()
    label(sl, f"{acc_pct:.1f}%", x+3.7+bw+0.08, y, 0.7, 0.25, size=10, color=tc)

def feat_bar(sl, feat, val, x, y):
    label(sl, feat, x, y, 1.8, 0.22, size=10, color=DARK)
    bw = val * 12.0
    s = sl.shapes.add_shape(1, Inches(x+1.9), Inches(y+0.06), Inches(bw), Inches(0.1))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
    label(sl, f"{val:.4f}", x+1.9+bw+0.08, y, 0.6, 0.22, size=9.5, color=MUTED)


# ── SLIDE 1: Problem Statement ─────────────────────────────────────────────────
sl1 = add_slide()
slide_num(sl1, 1)
title_bar(sl1, "Problem statement & objectives")

section_label(sl1, "Problem", 0.4, 1.15)
bullet_block(sl1, [
    "BTC price is highly volatile — next-day direction is difficult to anticipate from price alone",
    "Most retail participants have no systematic signal; decisions are largely reactive",
    "Standard statistical models struggle on non-stationary financial time series",
], 0.4, 1.4, 5.8)

section_label(sl1, "Objectives", 6.6, 1.15)
bullet_block(sl1, [
    "Predict next-day BTC/USD price direction (up or down) using ML",
    "Compare regression and classification models on the same dataset",
    "Build a live dashboard that translates model output into an actionable signal",
    "Identify which technical indicators carry the most predictive weight",
], 6.6, 1.4, 6.1)

tx = 0.4
ty = 5.9
for t, fill, tc in [("Binary classification", ACCENT2, ACCENT), ("Log returns", ACCENT2, ACCENT),
                    ("Time-series CV", ACCENT2, ACCENT), ("No look-ahead", AMBER2, AMBER)]:
    tag(sl1, t, tx, ty, fill, tc)
    tx += len(t)*0.085 + 0.35


# ── SLIDE 2: Dataset + EDA ──────────────────────────────────────────────────────
sl2 = add_slide()
slide_num(sl2, 2)
title_bar(sl2, "Dataset, EDA & preprocessing")

info_card(sl2, "Source", "1-minute OHLCV, BTC/USD\naggregated to daily bars", 0.4, 1.15, 3.9, 0.8)
info_card(sl2, "Train / test split", "80 / 20 (time-ordered)\n3,576 train · 895 test", 4.45, 1.15, 3.9, 0.8)
info_card(sl2, "Target variable", "y = sign(log return)\n1 if next-day return > 0, else 0", 8.5, 1.15, 4.3, 0.8)

section_label(sl2, "Preprocessing pipeline", 0.4, 2.2)

steps = ["Raw 1-min OHLCV", "Daily aggregation", "ADF stationarity", "Log returns", "20 features", "TimeSeriesSplit (5 folds)"]
sx = 0.4
for i, s in enumerate(steps):
    w = 1.85
    hl = i == len(steps)-1
    fill = ACCENT if hl else WHITE
    tc_  = WHITE if hl else DARK
    bx = box(sl2, sx, 2.45, w, 0.45, fill=fill)
    tf = sl2.shapes.add_textbox(Inches(sx+0.08), Inches(2.52), Inches(w-0.1), Inches(0.32))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s
    r.font.size = Pt(9.5); r.font.color.rgb = tc_; r.font.name = "Calibri"
    sx += w
    if i < len(steps)-1:
        label(sl2, "›", sx-0.05, 2.5, 0.15, 0.3, size=13, color=MUTED, align=PP_ALIGN.CENTER)

section_label(sl2, "Base indicators (10)", 0.4, 3.2)
tx = 0.4
for t in ["EMA 9/21", "MACD", "RSI 14", "MOM10", "PROC10", "StochK14", "OBV", "ATR"]:
    tag(sl2, t, tx, 3.45, ACCENT2, ACCENT)
    tx += len(t)*0.085 + 0.35

section_label(sl2, "Extended features (10)", 0.4, 3.95)
tx = 0.4
for t in ["Lag returns 1–5", "ATR norm", "BB %B", "52-week position", "Vol ratio", "Streak"]:
    tag(sl2, t, tx, 4.2, GREEN2, GREEN)
    tx += len(t)*0.085 + 0.35

section_label(sl2, "Key notes", 0.4, 5.0)
bullet_block(sl2, [
    "All features use shift(1) — strictly no look-ahead into future data",
    "ADF test confirmed log returns are stationary (p < 0.01)",
], 0.4, 5.25, 12.5)


# ── SLIDE 3: Methodology ────────────────────────────────────────────────────────
sl3 = add_slide()
slide_num(sl3, 3)
title_bar(sl3, "Methodology & models")

section_label(sl3, "Regression models", 0.4, 1.15)
bullet_block(sl3, [
    "Persistence baseline — naive carry-forward (lower bound)",
    "Linear regression, Ridge, Lasso — regularised OLS",
    "XGBoost regressor — gradient boosted trees, RMSE target, 8-param grid search",
], 0.4, 1.4, 5.9)

section_label(sl3, "Classification models", 0.4, 2.75)
bullet_block(sl3, [
    "Majority class baseline — always predicts the more frequent class",
    "Logistic regression — L2, C tuned via GridSearchCV",
    "Random forest — 100 estimators, limited depth",
    "XGBoost classifier — same grid as regressor, logloss eval",
    "Ensemble — XGBoost + Logistic averaged probability, majority vote",
], 0.4, 3.0, 5.9)

section_label(sl3, "Key design decisions", 6.8, 1.15)

for i, (t, sub) in enumerate([
    ("TimeSeriesSplit cross-validation", "5 folds, strictly forward — no leakage across time"),
    ("Threshold optimisation", "Decision boundary searched 0.40–0.65 on held-out 20% of training"),
    ("High-confidence filtering", "Signals reported separately when model probability exceeds 60%"),
    ("Ensemble averaging", "XGBoost + Logistic probabilities averaged before thresholding"),
]):
    info_card(sl3, t, sub, 6.8, 1.4 + i*1.05, 5.9, 0.95)


# ── SLIDE 4: Results ────────────────────────────────────────────────────────────
sl4 = add_slide()
slide_num(sl4, 4)
title_bar(sl4, "Results & evaluation")

section_label(sl4, "Classification accuracy", 0.4, 1.15)
models_cls = [
    ("Random forest",        50.2, False),
    ("Majority baseline",    51.2, False),
    ("XGBoost",              51.8, False),
    ("Logistic regression",  52.4, False),
    ("Ensemble (XGB + LR)",  53.5, True),
]
for i, (m, a, hl) in enumerate(models_cls):
    hbar(sl4, m, a, 0.4, 1.4 + i*0.42, highlight=hl)

section_label(sl4, "Regression — RMSE on log return", 6.8, 1.15)
reg_data = [
    ("Persistence baseline",  "0.03628", False),
    ("Linear regression",     "0.02537", False),
    ("XGBoost",               "0.02534", False),
    ("Ridge",                 "0.02509", False),
    ("Lasso",                 "0.02490", True),
]
for i, (m, v, hl) in enumerate(reg_data):
    tc = ACCENT if hl else DARK
    label(sl4, m, 6.8, 1.4 + i*0.42, 4.0, 0.3, size=10, color=tc)
    label(sl4, v, 11.0, 1.4 + i*0.42, 1.5, 0.3, size=10, color=tc)

metric_card(sl4, "53.5%", "Best accuracy — Ensemble",       0.4,  4.0, 3.0, vc=ACCENT)
metric_card(sl4, "59.9%", "Precision on UP predictions",    3.55, 4.0, 3.0)
metric_card(sl4, "51.0%", "High-conf accuracy (>60% prob)", 6.7,  4.0, 3.0)
metric_card(sl4, "20",    "Features used in best model",    9.85, 4.0, 3.0)

label(sl4,
    "BTC next-day direction is near-random — 53–54% is consistent with published ML results on efficient crypto markets. "
    "The Ensemble's 59.9% precision on UP calls is the most actionable metric for a directional strategy.",
    0.4, 5.05, 12.5, 0.5, size=10.5, color=MUTED)


# ── SLIDE 5: Demo + Conclusion ──────────────────────────────────────────────────
sl5 = add_slide()
slide_num(sl5, 5)
title_bar(sl5, "Demo, conclusions & future work")

section_label(sl5, "Dashboard built", 0.4, 1.15)
bullet_block(sl5, [
    "React + Recharts frontend — live price chart with gradient area fill",
    "Split-screen: price chart left, ML signal panel right (UP / DOWN with confidence)",
    "Composite signal −7 to +7 from RSI, MACD, EMA21, momentum, ML probability",
    "Flask REST API — 7 endpoints: metrics, predictions, feature importance, advisor",
], 0.4, 1.4, 5.9)

section_label(sl5, "Conclusions", 0.4, 3.45)
bullet_block(sl5, [
    "No single model dominates — ensemble averaging reduces variance",
    "Lag returns and ATR are the top two predictors; RSI and BB %B follow",
    "Market efficiency caps accuracy; high-confidence filtering is a practical workaround",
], 0.4, 3.7, 5.9)

section_label(sl5, "Top features by importance", 6.8, 1.15)
feats = [
    ("lag_return_1", 0.0646),
    ("atr_norm",     0.0599),
    ("RSI",          0.0583),
    ("bb_pct_b",     0.0580),
    ("streak",       0.0570),
    ("Volume",       0.0551),
    ("lag_return_2", 0.0534),
]
for i, (f, v) in enumerate(feats):
    feat_bar(sl5, f, v, 6.8, 1.4 + i*0.4)

section_label(sl5, "Future work", 6.8, 4.5)
bullet_block(sl5, [
    "Incorporate on-chain data — active addresses, exchange flows",
    "Sentiment features from news headlines or social volume",
    "Walk-forward retraining on a rolling window",
    "Backtested P&L to validate signal quality beyond accuracy",
], 6.8, 4.75, 5.9)


prs.save(r"D:\ML Proj\BTC_ML_Presentation.pptx")
print("Saved: D:\\ML Proj\\BTC_ML_Presentation.pptx")
